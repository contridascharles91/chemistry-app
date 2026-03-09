"""
Chunks Chemistry - Production Server
Cloud-ready with R2 storage integration
"""

from flask import Flask, request, jsonify, Response, stream_with_context, abort
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from werkzeug.utils import secure_filename
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import json
import os
import re
import math
import logging
from datetime import datetime
from collections import Counter
import PyPDF2
import docx
from pptx import Presentation

# ── Logging ───────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

# ── HTTP Session — connection pooling + retry ─────────────────────────────────
def _build_session():
    session = requests.Session()
    retry = Retry(total=3, backoff_factor=0.5,
                  status_forcelist=[502, 503, 504],
                  allowed_methods=["GET", "POST"], raise_on_status=False)
    adapter = HTTPAdapter(max_retries=retry, pool_connections=10, pool_maxsize=20)
    session.mount("https://", adapter)
    session.mount("http://",  adapter)
    return session

_session = _build_session()

try:
    import numpy as np
    NUMPY_AVAILABLE = True
except ImportError:
    NUMPY_AVAILABLE = False
    logger.warning("numpy not installed — semantic search disabled")

# ── Allowed file types ────────────────────────────────────────────────────────
ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.ppt'}

def allowed_file(filename):
    return os.path.splitext(filename.lower())[1] in ALLOWED_EXTENSIONS

def sanitize_text(text, max_len=2000):
    text = str(text).replace('\x00', '').strip()
    return text[:max_len]

app = Flask(__name__)

# ── CORS ──────────────────────────────────────────────────────────────────────
# FIX 8: Build allowed origins from environment so staging/preview URLs work.
# FRONTEND_URL can be a comma-separated list of origins via ALLOWED_ORIGINS env var.
FRONTEND_URL = os.environ.get('FRONTEND_URL', 'http://localhost:5173')
BACKEND_URL  = os.environ.get('BACKEND_URL',  'http://localhost:5000')

_raw_origins = os.environ.get('ALLOWED_ORIGINS', '')
_allowed_origins = [o.strip() for o in _raw_origins.split(',') if o.strip()] if _raw_origins else []
_default_origins = [
    "https://chunks-ai.vercel.app",
    "http://localhost:5173",
    "http://localhost:3000",
    "http://localhost:5000",
]
if FRONTEND_URL and FRONTEND_URL not in _default_origins:
    _default_origins.append(FRONTEND_URL)
CORS_ORIGINS = list(dict.fromkeys(_allowed_origins + _default_origins))

# FIX: Also allow all *.vercel.app subdomains (covers preview deploys and admin.html)
# flask-cors 4.x supports re.compile() objects mixed into the origins list
import re as _re
CORS_ORIGINS.append(_re.compile(r'^https://[a-zA-Z0-9-]+\.vercel\.app$'))

CORS(app,
     origins=CORS_ORIGINS,
     allow_headers=["Content-Type", "Authorization", "X-Requested-With"],
     methods=["GET", "POST", "OPTIONS"],
     supports_credentials=False,
     max_age=86400)


@app.after_request
def after_request(response):
    # Security headers — CORS handled by flask-cors above
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'DENY'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Referrer-Policy'] = 'strict-origin-when-cross-origin'
    response.headers['Permissions-Policy'] = 'camera=(), microphone=(), geolocation=()'
    response.headers.pop('Server', None)
    response.headers.pop('X-Powered-By', None)
    return response


# ── Rate Limiting ─────────────────────────────────────────────────────────────
limiter = Limiter(
    key_func=get_remote_address,
    app=app,
    default_limits=["200 per hour", "30 per minute"],
    storage_uri="memory://",
    strategy="fixed-window"
)

# ── Upload size limit 25MB ────────────────────────────────────────────────────
app.config['MAX_CONTENT_LENGTH'] = 25 * 1024 * 1024
# ============================================

OPENROUTER_API_KEY = os.environ.get('OPENROUTER_API_KEY', 'your-key-here')
R2_BUCKET_URL      = os.environ.get('R2_BUCKET_URL', 'https://pub-xxxxx.r2.dev')
# FIX: warn loudly at startup if placeholder values are still in use
if R2_BUCKET_URL == 'https://pub-xxxxx.r2.dev':
    logger.warning("⚠️  R2_BUCKET_URL is still the placeholder value — all book URLs will 404! Set R2_BUCKET_URL in your environment.")
if OPENROUTER_API_KEY == 'your-key-here':
    logger.warning("⚠️  OPENROUTER_API_KEY is not set — all AI calls will fail! Set OPENROUTER_API_KEY in your environment.")
PORT               = int(os.environ.get('PORT', 5000))
PRODUCTION         = os.environ.get('PRODUCTION', 'false').lower() == 'true'
OPENROUTER_URL     = "https://openrouter.ai/api/v1/chat/completions"
MODEL              = os.environ.get('MODEL', 'arcee-ai/trinity-large-preview:free')

# Supabase config for server-side JWT verification and daily message counters
SUPABASE_URL         = os.environ.get('SUPABASE_URL', '')
SUPABASE_SERVICE_KEY = os.environ.get('SUPABASE_SERVICE_KEY', '')  # service_role key (secret)
SUPABASE_ANON_KEY    = os.environ.get('SUPABASE_ANON_KEY', '')     # anon/public key (safe to expose, but not in source)

# In-memory free-tier message counter (fallback: keyed by IP+date when Supabase is unavailable)
_free_tier_counters = {}

FREE_TIER_DAILY_LIMIT = 20   # matches the 20-message client-side limit
MAX_HISTORY_TURNS     = 10   # consistent conversation context window across all AI callers


def _verify_supabase_jwt(token: str) -> dict | None:
    """
    Verify a Supabase JWT and return the user record from the DB, or None if invalid.
    Uses the Supabase REST API so no extra Python libraries are needed.
    Returns dict with at least: {'id': <uuid>, 'email': <str>}
    """
    if not SUPABASE_URL or not SUPABASE_SERVICE_KEY or not token:
        return None
    try:
        resp = _session.get(
            f"{SUPABASE_URL}/auth/v1/user",
            headers={
                "Authorization": f"Bearer {token}",
                "apikey": SUPABASE_SERVICE_KEY,
            },
            timeout=5
        )
        if resp.status_code == 200:
            return resp.json()
    except Exception as e:
        logger.warning(f"JWT verify error: {e}")
    return None


def _get_user_tier_from_db(user_id: str) -> str:
    """
    Look up the user's subscription tier in Supabase.
    Returns 'paid' or 'free'. Defaults to 'free' on any error.
    Expects a 'users' table with columns: id (uuid), tier (text).
    """
    if not SUPABASE_URL or not SUPABASE_SERVICE_KEY or not user_id:
        return 'free'
    try:
        resp = _session.get(
            f"{SUPABASE_URL}/rest/v1/users",
            params={"id": f"eq.{user_id}", "select": "tier"},
            headers={
                "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
                "apikey": SUPABASE_SERVICE_KEY,
            },
            timeout=5
        )
        if resp.status_code == 200:
            rows = resp.json()
            if rows and rows[0].get('tier') in ('paid', 'pro', 'ultra'):
                return 'paid'
    except Exception as e:
        logger.warning(f"Tier lookup error: {e}")
    return 'free'


def _get_and_increment_daily_count(user_id: str, date_str: str) -> int:
    """
    Atomically read + increment a daily message counter for a user in Supabase.
    Table: free_tier_usage (user_id uuid, date date, count int)
    Returns the NEW count after incrementing.
    Falls back to in-memory counter if Supabase is unavailable.
    """
    if not SUPABASE_URL or not SUPABASE_SERVICE_KEY:
        # Fallback: in-memory counter keyed by user_id (or IP) + date
        key = f"freetier:{user_id}:{date_str}"
        count = _free_tier_counters.get(key, 0) + 1
        _free_tier_counters[key] = count
        return count

    try:
        # Upsert: increment count or insert 1 if no row yet.
        # Uses Supabase RPC for atomic increment to avoid race conditions.
        resp = _session.post(
            f"{SUPABASE_URL}/rest/v1/rpc/increment_free_tier_usage",
            json={"p_user_id": user_id, "p_date": date_str},
            headers={
                "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
                "apikey": SUPABASE_SERVICE_KEY,
                "Content-Type": "application/json",
            },
            timeout=5
        )
        if resp.status_code == 200:
            result = resp.json()
            # RPC returns the new count as an integer
            if isinstance(result, int):
                return result
            if isinstance(result, dict):
                return result.get('count', 1)
    except Exception as e:
        logger.warning(f"Supabase daily count error: {e}")

    # Fallback to in-memory on any Supabase failure
    key = f"freetier:{user_id}:{date_str}"
    count = _free_tier_counters.get(key, 0) + 1
    _free_tier_counters[key] = count
    return count


# ============================================
# BOOK LIBRARY - R2 URLs
# ============================================

BOOK_LIBRARY = {
    'zumdahl': {
        'name': 'General Chemistry',
        'author': 'Zumdahl & Zumdahl',
        'chunks_url': f'{R2_BUCKET_URL}/data/zumdhal_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/zumdhal.pdf'
    },
    'atkins': {
        'name': 'Physical Chemistry',
        'author': 'Atkins & de Paula',
        'chunks_url': f'{R2_BUCKET_URL}/data/atkins_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/atkins_physical_chemistry.pdf'
    },
    'harris': {
        'name': 'Quantitative Chemical Analysis',
        'author': 'Daniel C. Harris',
        'chunks_url': f'{R2_BUCKET_URL}/data/harris_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/harris_quantitative_analysis.pdf'
    },
    'klein': {
        'name': 'Organic Chemistry',
        'author': 'David Klein',
        'chunks_url': f'{R2_BUCKET_URL}/data/klein_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/klein_organic_chemistry.pdf'
    },
    'berg': {
        'name': 'Biochemistry',
        'author': 'Berg, Tymoczko & Stryer',
        'chunks_url': f'{R2_BUCKET_URL}/data/berg_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/berg_biochemistry.pdf'
    },    
    'netter': {
        'name': 'Atlas of Human Anatomy',
        'author': 'Frank H. Netter',
        'chunks_url': f'{R2_BUCKET_URL}/data/atlas_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/Atlas.pdf'
    },
    'anaphy2e': {
        'name': 'Anatomy & Physiology',
        'author': 'Patton & Thibodeau',
        'chunks_url': f'{R2_BUCKET_URL}/data/anaphy2e_chunks_with_embeddings.json',
        'pdf_url':    f'{R2_BUCKET_URL}/data/anaphy2e.pdf'
    }
}
# FIX: Backward-compat alias so old frontend bookId='biochemistry' still works
BOOK_LIBRARY['biochemistry'] = BOOK_LIBRARY['berg']


# ============================================
# KEYWORD-BASED TEXTBOOK SEARCH (TF-IDF)
# ============================================

STOPWORDS = {
    'a', 'an', 'the', 'is', 'it', 'in', 'on', 'at', 'to', 'for',
    'of', 'and', 'or', 'but', 'with', 'this', 'that', 'are', 'was',
    'be', 'as', 'by', 'from', 'what', 'how', 'why', 'when', 'where',
    'which', 'who', 'do', 'does', 'did', 'can', 'could', 'would',
    'should', 'will', 'have', 'has', 'had', 'not', 'if', 'so', 'its'
}

def tokenize(text):
    words = re.findall(r'[a-z]+', text.lower())
    return [w for w in words if w not in STOPWORDS and len(w) > 2]

def tfidf_score(query_tokens, chunk_tokens, idf_map):
    tf = Counter(chunk_tokens)
    total = len(chunk_tokens) or 1
    score = 0.0
    for token in query_tokens:
        tf_val = tf.get(token, 0) / total
        idf_val = idf_map.get(token, 0)
        score += tf_val * idf_val
    return score

def enhanced_score(query_tokens, chunk, chunk_tokens, idf_map):
    base = tfidf_score(query_tokens, chunk_tokens, idf_map)
    text_lower = chunk.get('text', '').lower()
    total_words = len(chunk_tokens) or 1
    bonus = 0.0

    first_100 = text_lower[:100]
    for token in query_tokens:
        if text_lower.startswith(token) or f'{token} is' in first_100 or f'{token} are' in first_100:
            bonus += 0.015
        first_20pct = text_lower[:max(50, len(text_lower) // 5)]
        if token in first_20pct:
            bonus += 0.008

    query_phrase = ' '.join(query_tokens)
    if query_phrase in text_lower:
        bonus += 0.02
    for i in range(len(query_tokens) - 1):
        bigram = query_tokens[i] + ' ' + query_tokens[i + 1]
        if bigram in text_lower:
            bonus += 0.008

    tf = Counter(chunk_tokens)
    query_word_count = sum(tf.get(t, 0) for t in query_tokens)
    density = query_word_count / total_words
    bonus += min(density * 0.5, 0.025)

    if total_words > 300 and density < 0.02:
        bonus -= 0.005
    if total_words < 150 and base > 0:
        bonus += 0.005

    return base + bonus


# FIX: Module-level query embedding cache shared across all book instances.
# Keyed by question text — embedding only depends on text, not which book is loaded.
_global_query_cache: dict = {}

class TextbookSearch:
    # Cosine similarity threshold for hybrid search (embeddings active):
    # 0.0–1.0 scale. 0.25 = weak but plausible match. Below → answer from general knowledge.
    LOW_CONFIDENCE_HYBRID  = 0.25
    # TF-IDF-only threshold (fallback when embeddings unavailable):
    LOW_CONFIDENCE_TFIDF   = 0.010

    EMBEDDING_MODEL = "openai/text-embedding-3-small"
    EMBEDDING_DIMS  = 1536

    def __init__(self):
        self.chunks            = []
        self.tokenized_chunks  = []
        self.idf_map           = {}
        self.book_id           = None
        self.embedding_matrix  = None   # numpy array shape (N, 1536) — None until loaded
        self.has_embeddings    = False
        self._query_cache      = _global_query_cache  # FIX: shared module-level cache, not per-instance

    def load_chunks_from_url(self, url, book_id=None):
        try:
            logger.info(f"📥 Fetching chunks from: {url}")
            response = _session.get(url, timeout=60)
            response.raise_for_status()
            chunks = response.json()
            logger.info(f"✅ Loaded {len(chunks)} chunks")

            self.chunks   = chunks
            self.book_id  = book_id

            # ── Build TF-IDF index (always, used as fallback / hybrid component) ──
            self.tokenized_chunks = [tokenize(c.get('text', '')) for c in chunks]
            N  = len(self.tokenized_chunks)
            df = Counter()
            for tokens in self.tokenized_chunks:
                for t in set(tokens):
                    df[t] += 1
            self.idf_map = {t: math.log((N + 1) / (df[t] + 1)) for t in df}
            logger.info(f"✅ TF-IDF index: {len(self.idf_map)} unique terms")

            # ── Load embedding matrix if embeddings are present ──────────────────
            try:
                import numpy as np
                first_with_emb = next((c for c in chunks if c.get('embedding')), None)
                if first_with_emb:
                    dims = len(first_with_emb['embedding'])
                    if dims == self.EMBEDDING_DIMS:
                        # FIX: float32 for better cosine similarity precision (float16 loses too much precision)
                        matrix = np.array(
                            [c['embedding'] for c in chunks],
                            dtype=np.float32
                        )
                        # L2-normalise rows so dot product == cosine similarity
                        norms = np.linalg.norm(matrix, axis=1, keepdims=True)
                        norms = np.where(norms == 0, 1, norms)   # avoid div-by-zero
                        self.embedding_matrix = matrix / norms
                        self.has_embeddings   = True
                        logger.info(f"✅ Embedding matrix loaded: {matrix.shape}")
                    else:
                        logger.warning(f"⚠️  Embedding dims={dims}, expected {self.EMBEDDING_DIMS}. "
                              f"Re-run process_book.py to regenerate. Falling back to TF-IDF.")
                else:
                    logger.info("No embeddings in JSON — TF-IDF only."
                          "Run process_book.py to add semantic search.")
            except ImportError:
                logger.warning("numpy not installed — semantic search disabled. Run: pip install numpy")

            return True

        except Exception as e:
            logger.error(f"Error loading chunks: {e}")
            logger.exception("Unhandled error")
            return False

    # ── Query embedding ───────────────────────────────────────────────────────

    def _embed_query(self, text: str) -> 'np.ndarray | None':
        """
        Embed a query string using OpenRouter. Returns a normalised float32
        vector of shape (1536,), or None if the call fails.

        Results are cached in _query_cache so repeated identical questions
        (e.g. exam regeneration) never make a second API call.
        """
        if text in self._query_cache:
            return self._query_cache[text]

        try:
            import numpy as np
            headers = {
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "Content-Type":  "application/json",
                "HTTP-Referer":  "https://chunks-ai.vercel.app",
                "X-Title":       "Chunks Chemistry"
            }
            payload = {
                "model": self.EMBEDDING_MODEL,
                "input": [text]
            }
            resp = _session.post(
                "https://openrouter.ai/api/v1/embeddings",
                headers=headers, json=payload, timeout=15
            )
            if resp.status_code == 200:
                vec  = np.array(resp.json()["data"][0]["embedding"], dtype=np.float32)  # FIX: float32 for precision
                norm = np.linalg.norm(vec)
                if norm > 0:
                    vec /= norm
                self._query_cache[text] = vec
                return vec
            logger.warning(f"Embedding API {resp.status_code}: {resp.text[:200]}")
        except Exception as e:
            logger.warning(f"Embedding query failed: {e}")

        return None

    # ── Search ────────────────────────────────────────────────────────────────

    def smart_search(self, question: str, top_k: int = 5):
        """
        Return (context_str, score, is_relevant, best_source, all_sources).

        Scoring strategy:
        ┌─────────────────────────────────────────────────────┐
        │  embeddings available?                              │
        │    YES → hybrid score = 0.70 × cosine              │
        │                        + 0.30 × tfidf_normalised   │
        │    NO  → tfidf score only (fallback)                │
        └─────────────────────────────────────────────────────┘

        Cosine catches semantic similarity ("why does ice float" ↔
        "hydrogen bonding anomalous density"). TF-IDF reinforces exact
        term matches (formula names, specific compound names).
        """
        if not self.chunks:
            return "No textbook loaded.", 0.0, False, None, []

        query_tokens = tokenize(question)

        # ── 1. TF-IDF scores (always computed) ───────────────────────────────
        tfidf_scores = []
        for i, chunk in enumerate(self.chunks):
            s = enhanced_score(query_tokens, chunk, self.tokenized_chunks[i], self.idf_map) \
                if query_tokens else 0.0
            tfidf_scores.append(s)

        # Normalise TF-IDF to [0, 1] so it's on the same scale as cosine
        max_tfidf = max(tfidf_scores) if tfidf_scores else 1.0
        if max_tfidf > 0:
            tfidf_norm = [s / max_tfidf for s in tfidf_scores]
        else:
            tfidf_norm = tfidf_scores

        # ── 2. Cosine similarity (when embeddings loaded) ─────────────────────
        use_hybrid    = False
        final_scores  = list(tfidf_norm)   # default: tfidf only
        low_conf      = self.LOW_CONFIDENCE_TFIDF

        if self.has_embeddings:
            import numpy as np
            query_vec = self._embed_query(question)
            if query_vec is not None:
                # Matrix dot product: shape (N,) of cosine similarities
                cosine = self.embedding_matrix.dot(query_vec)
                # Clip to [0, 1] — negative cosine for textbook search = irrelevant
                cosine = cosine.clip(0, 1).tolist()

                # Hybrid: 70% semantic + 30% keyword
                final_scores = [
                    0.70 * cos + 0.30 * tfidf
                    for cos, tfidf in zip(cosine, tfidf_norm)
                ]
                use_hybrid = True
                low_conf   = self.LOW_CONFIDENCE_HYBRID
            else:
                logger.warning("Query embedding failed — falling back to TF-IDF")

        # ── 3. Rank and select ────────────────────────────────────────────────
        scored = sorted(
            zip(self.chunks, final_scores),
            key=lambda x: x[1], reverse=True
        )

        top_score  = scored[0][1] if scored else 0.0
        is_relevant = top_score >= low_conf

        mode_label = "hybrid" if use_hybrid else "tfidf"
        logger.debug(f"[{mode_label}] score={top_score:.4f} relevant={is_relevant}")

        context = "\n\n".join([
            f"[Page {c['page']}] {c['text']}"
            for c, _ in scored[:top_k]
        ])

        all_sources = [
            {'page': int(c['page']), 'text': c.get('text', '')[:200]}
            for c, s in scored[:top_k]
            if s >= low_conf
        ]

        best_source = all_sources[0] if all_sources else None
        return context, top_score, is_relevant, best_source, all_sources

    def get_candidate_pages(self, topic: str, top_k: int = 5):
        if not self.chunks:
            return []
        query_tokens = tokenize(topic)
        scored = []
        for i, chunk in enumerate(self.chunks):
            score = enhanced_score(query_tokens, chunk, self.tokenized_chunks[i], self.idf_map)
            scored.append({'page': chunk['page'], 'text': chunk['text'], 'score': score})
        scored.sort(key=lambda x: x['score'], reverse=True)
        return scored[:top_k]


# ── In-memory material cache — pre-generated flashcards & quizzes ────────────
# Key: (book_id, topic_normalised, material_type, count)  Value: (result, timestamp)
import time as _time
_material_cache: dict = {}
_MATERIAL_CACHE_TTL = 86400  # 24 hours — regenerate daily

def _cache_key(book_id: str, topic: str, mtype: str, count: int) -> str:
    norm = re.sub(r'[^a-z0-9]', '_', topic.lower().strip())[:60]
    return f"{mtype}:{book_id}:{norm}:{count}"

def _cache_get(key: str):
    entry = _material_cache.get(key)
    if not entry:
        return None
    result, ts = entry
    if _time.time() - ts > _MATERIAL_CACHE_TTL:
        del _material_cache[key]
        return None
    return result

def _cache_set(key: str, value) -> None:
    # Limit cache to 500 entries — evict oldest on overflow
    if len(_material_cache) >= 500:
        oldest = min(_material_cache.items(), key=lambda x: x[1][1])
        del _material_cache[oldest[0]]
    _material_cache[key] = (value, _time.time())

# Per-book index cache — each book_id maps to its own TextbookSearch instance.
# Eliminates multi-worker race conditions where two users on different books
# would overwrite the single global index.
_book_cache: dict[str, TextbookSearch] = {}

def get_book_index(book_id: str) -> TextbookSearch:
    """Return a cached (or freshly loaded) TextbookSearch for book_id."""
    if book_id in _book_cache:
        return _book_cache[book_id]

    if book_id not in BOOK_LIBRARY:
        return TextbookSearch()  # empty — is_relevant will be False

    searcher = TextbookSearch()
    ok = searcher.load_chunks_from_url(BOOK_LIBRARY[book_id]['chunks_url'], book_id=book_id)
    if ok:
        _book_cache[book_id] = searcher
        mode = "hybrid (embeddings + TF-IDF)" if searcher.has_embeddings else "TF-IDF only"
        logger.info(f"✅ Cached [{mode}] index for: {book_id}")
    return searcher


# ============================================
# AI CALLER
# ============================================

def call_ai(prompt, system_prompt="You are an expert chemistry tutor.", model=None,
            history=None, max_tokens_override=None):
    use_model = model or MODEL
    try:
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://chunks-ai.vercel.app",
            "X-Title": "Chunks Chemistry"
        }
        messages = [{"role": "system", "content": system_prompt}]
        if history:
            for h in history[-MAX_HISTORY_TURNS:]:
                role = h.get("role", "user")
                content = h.get("content", "")
                if role in ("user", "assistant") and content:
                    messages.append({"role": role, "content": content})
        messages.append({"role": "user", "content": prompt})

        payload = {
            "model": use_model,
            "messages": messages,
            # FIX: lowered temperature from 0.4 → 0.15
            # Chemistry facts, equations, and constants must be deterministic.
            "temperature": 0.15,
            "max_tokens": max_tokens_override if max_tokens_override else 6000
        }
        logger.info(f"Model: {use_model} | history: {len(history) if history else 0} turns")
        response = _session.post(OPENROUTER_URL, headers=headers, json=payload, timeout=55)
        if response.status_code == 200:
            resp_json = response.json()
            choices = resp_json.get('choices', [])
            if choices:
                return choices[0]['message']['content']
            err = resp_json.get('error', {})
            return f"Error: Model returned no choices — {err.get('message', str(resp_json)[:200])}"
        logger.error(f"API error {response.status_code}: {response.text[:300]}")
        return f"Error: API returned {response.status_code} - {response.text[:200]}"
    except requests.Timeout:
        return "Error: The AI model timed out. Please try again."
    except Exception as e:
        logger.exception("Unhandled error")
        return f"Error: {str(e)}"


def call_ai_web_search(question, system_prompt=None, history=None):
    """
    Uses Perplexity Sonar via OpenRouter for real-time web search with citations.
    Returns (answer_text, citations_list)
    citations_list is a list of dicts: [{url, title}]
    """
    WEB_MODEL = os.environ.get('WEB_MODEL', 'perplexity/sonar')
    try:
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type":  "application/json",
            "HTTP-Referer":  "https://chunks-ai.vercel.app",
            "X-Title":       "Chunks Chemistry"
        }

        sys_prompt = system_prompt or (
            "You are a helpful research assistant. Answer clearly and accurately using "
            "current web information. Always include specific references to the sources "
            "you used. Format your answer in clean markdown with headers where appropriate."
        )

        messages = [{"role": "system", "content": sys_prompt}]
        if history:
            for h in (history or [])[-MAX_HISTORY_TURNS:]:  # FIX: was -6, now consistent with call_ai
                role    = h.get("role", "user")
                content = h.get("content", "")
                if role in ("user", "assistant") and content:
                    messages.append({"role": role, "content": content})
        messages.append({"role": "user", "content": question})

        payload = {
            "model": WEB_MODEL,
            "messages": messages,
            "temperature": 0.2,
            "max_tokens": 4000
        }

        logger.info(f"Web search model: {WEB_MODEL} | Q: {question[:80]}")
        response = _session.post(OPENROUTER_URL, headers=headers, json=payload, timeout=60)

        if response.status_code != 200:
            logger.error(f"Web search API error {response.status_code}: {response.text[:300]}")
            return f"Web search error: {response.status_code}", []

        resp_json = response.json()
        choices   = resp_json.get('choices', [])
        if not choices:
            return "No results returned.", []

        answer = choices[0]['message']['content']

        # Perplexity via OpenRouter returns citations at top level or in choices
        raw_citations = (
            resp_json.get('citations') or
            choices[0].get('message', {}).get('citations') or
            choices[0].get('delta', {}).get('citations') or
            []
        )

        # Normalize citations — can be strings (URLs) or dicts
        citations = []
        seen_urls = set()
        for c in raw_citations:
            if isinstance(c, str) and c.startswith('http'):
                url = c
                # Try to derive a title from domain
                try:
                    from urllib.parse import urlparse
                    domain = urlparse(url).netloc.replace('www.', '')
                    title  = domain
                except Exception:
                    title = url
                if url not in seen_urls:
                    seen_urls.add(url)
                    citations.append({'url': url, 'title': title})
            elif isinstance(c, dict):
                url   = c.get('url', '')
                title = c.get('title') or c.get('name') or ''
                if not title and url:
                    try:
                        from urllib.parse import urlparse
                        title = urlparse(url).netloc.replace('www.', '')
                    except Exception:
                        title = url
                if url and url not in seen_urls:
                    seen_urls.add(url)
                    citations.append({'url': url, 'title': title})

        # Also extract any URLs embedded in the answer text as fallback
        if not citations:
            found_urls = re.findall(r'https?://[^\s\)\]\>\"\']+', answer)
            for url in found_urls:
                url = url.rstrip('.,;:')
                if url not in seen_urls:
                    seen_urls.add(url)
                    try:
                        from urllib.parse import urlparse
                        title = urlparse(url).netloc.replace('www.', '')
                    except Exception:
                        title = url
                    citations.append({'url': url, 'title': title})

        logger.info(f"Web search complete | citations: {len(citations)}")
        return answer, citations

    except requests.Timeout:
        return "Error: Web search timed out. Please try again.", []
    except Exception as e:
        logger.exception("Web search error")
        return f"Error: {str(e)}", []


# ============================================
# OFF-TOPIC DETECTION
# ============================================

# ── Subject-agnostic textbook search gating ──────────────────────────────────
#
# OLD APPROACH (broken): hardcoded subject keyword list — any book outside
# chemistry (nursing, biology, physics, etc.) would always skip textbook search
# because none of its domain words were in the list.
#
# NEW APPROACH: if a book is loaded, ALWAYS attempt textbook search.
# The LOW_CONFIDENCE threshold (0.010) already filters out bad matches —
# if the retrieved chunks aren't relevant, is_relevant=False and the AI
# answers from general knowledge anyway. No need to pre-filter by subject.
#
# The only questions we skip are obvious non-study chit-chat where searching
# a textbook would never help regardless of subject.

# Patterns that are clearly not study questions — short social exchanges,
# UI commands, personal questions about the AI itself.
_SKIP_PATTERNS = re.compile(
    r'^(hi+|hey+|hello|howdy|sup|yo+|hiya)[!?,.\s]*$'           # pure greetings
    r'|^(thanks?|thank you|thx|ty|tysm)[!?,.\s]*$'              # thank-yous
    r'|^(ok|okay|got it|sure|cool|nice|great|perfect)[!?,.\s]*$'# one-word acks
    r'|^(who (are|made|created|built) you)'                      # AI identity
    r'|^(what (is your name|can you do|are you))'               # AI capability
    r'|^(how are you|are you (ok|good|alive|sentient))'         # AI wellbeing
    r'|^(lol|lmao|haha|hehe|😂|👍|🙏)[!?,.\s]*$',              # reactions
    re.IGNORECASE
)

def should_search_textbook(question: str, chunks_loaded: bool) -> bool:
    """
    Return True if we should search the loaded textbook for this question.

    Rules:
    - If no book is loaded → False (nothing to search)
    - If the question is obvious non-study chit-chat → False
    - Everything else → True (let the LOW_CONFIDENCE threshold decide relevance)

    This replaces the old subject-keyword approach which blocked textbook search
    for any non-chemistry subject (nursing, biology, physics, etc.).
    """
    if not chunks_loaded:
        return False
    q = question.strip()
    if not q:
        return False
    # Skip very short non-questions (< 3 words and no '?')
    if len(q.split()) < 3 and '?' not in q and not any(c.isdigit() for c in q):
        if _SKIP_PATTERNS.match(q):
            return False
    return True


# ============================================
# ROUTES
# ============================================


# ============================================
# ERROR HANDLERS
# ============================================

@app.errorhandler(413)
def too_large(e):
    return jsonify({'success': False, 'error': 'File too large. Maximum is 25 MB.'}), 413

@app.errorhandler(429)
def rate_limited(e):
    return jsonify({'success': False, 'error': 'Too many requests. Please slow down.'}), 429

@app.errorhandler(404)
def not_found(e):
    return jsonify({'success': False, 'error': 'Endpoint not found.'}), 404

@app.errorhandler(405)
def method_not_allowed(e):
    return jsonify({'success': False, 'error': 'Method not allowed.'}), 405

@app.errorhandler(500)
def internal_error(e):
    logger.exception("Unhandled 500 error")
    return jsonify({'success': False, 'error': 'Internal server error.'}), 500

@app.route('/')
def home():
    return jsonify({
        'name': 'Chunks Chemistry API',
        'version': '2.0',
        'status': 'running',
        'endpoints': {
            'health': '/health',
            'ask': '/ask',
            'load_book': '/load-book',
            'pdf': '/pdf/<book_id>',
            'library': '/get-library',
            'flashcards': '/generate-flashcards',
            'upload_document': '/upload-document',
            'study_materials': '/generate-study-materials',
            'quiz': '/generate-quiz',
            'ask_image': '/ask-image'
        }
    })


@app.route('/ping', methods=['GET'])
def ping():
    result = call_ai("Reply with only the word: OK", system_prompt="You are a test bot.", model=MODEL)
    return jsonify({
        'success': 'Error' not in result,
        'model': MODEL,
        'api_key_set': OPENROUTER_API_KEY != 'your-key-here',
        'ai_response': result
    })


# FIX: was missing @app.route decorator — /health was a 404
@app.route('/api/config', methods=['GET', 'OPTIONS'])
@limiter.limit('30 per minute; 200 per hour')  # FIX: add rate limit to prevent key enumeration
def get_client_config():
    """Return public config values the frontend needs (no secrets here)."""
    return jsonify({
        'supabaseUrl':     SUPABASE_URL,
        'supabaseAnonKey': SUPABASE_ANON_KEY,
    })


@app.route('/health', methods=['GET'])
def health():
    book_status = {}
    for bid, searcher in _book_cache.items():
        book_status[bid] = {
            'chunks': len(searcher.chunks),
            'search_mode': 'hybrid' if searcher.has_embeddings else 'tfidf'
        }
    return jsonify({
        'status': 'healthy',
        'mode': 'production' if PRODUCTION else 'development',
        'books_cached': book_status,
        'books_available': list(BOOK_LIBRARY.keys()),
        'r2_configured': R2_BUCKET_URL != 'https://pub-xxxxx.r2.dev',
        'api_configured': OPENROUTER_API_KEY != 'your-key-here'
    })


@app.route('/load-book', methods=['POST', 'OPTIONS'])
@limiter.limit('10 per minute; 30 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def load_book():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        data = request.json
        book_id = data.get('bookId')
        logger.info(f"Load book request: {book_id}")

        if book_id not in BOOK_LIBRARY:
            return jsonify({'success': False, 'error': f'Book "{book_id}" not found'}), 404

        book = BOOK_LIBRARY[book_id]
        searcher = get_book_index(book_id)

        if not searcher.chunks:
            return jsonify({'success': False, 'error': 'Failed to load chunks from R2'}), 500

        return jsonify({
            'success': True,
            'book_id': book_id,
            'book_name': book['name'],
            'author': book['author'],
            'chunks_count': len(searcher.chunks)
        })

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/pdf/<book_id>', methods=['GET'])
@limiter.limit('20 per minute; 100 per hour')
def serve_pdf(book_id):
    if book_id not in BOOK_LIBRARY:
        return jsonify({'error': 'Book not found'}), 404

    pdf_url = BOOK_LIBRARY[book_id]['pdf_url']
    logger.info(f"Proxying PDF for: {book_id}")
    try:
        r = requests.get(pdf_url, timeout=60, stream=True)
        r.raise_for_status()
        return Response(
            stream_with_context(r.iter_content(chunk_size=8192)),
            content_type='application/pdf',
            headers={
                'Content-Disposition': f'inline; filename="{book_id}.pdf"',
                'Access-Control-Allow-Origin': '*'
            }
        )
    except Exception as e:
        logger.error(f"PDF proxy error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/get-library', methods=['GET'])
def get_library():
    books = [
        {'id': bid, 'name': info['name'], 'author': info['author'], 'available': True}
        for bid, info in BOOK_LIBRARY.items()
    ]
    return jsonify({'success': True, 'books': books})


# ============================================
# /ask — MAIN CHAT ENDPOINT
# ============================================

@app.route('/ask', methods=['POST', 'OPTIONS'])
@limiter.limit('20 per minute; 100 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def ask():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        data = request.json
        if not data:
            return jsonify({'success': False, 'error': 'Invalid or missing JSON body'}), 400

        question      = data.get('question', '')
        complexity    = max(1, min(10, int(data.get('complexity', 3))))
        mode          = data.get('mode', 'study').lower().strip()
        book_id       = data.get('bookId', 'zumdahl')
        thinking_mode = data.get('thinking', None)
        web_search    = data.get('web_search', False)
        history       = data.get('history', [])
        user_memory   = data.get('user_memory', '')

        # ── Server-side tier verification (fixes Blocker 2) ──────────────────
        # NEVER trust the client-sent user_tier — always verify via Supabase JWT.
        # Flow:
        #   1. Extract Bearer token from Authorization header
        #   2. Verify token with Supabase → get user_id
        #   3. Look up tier in DB (not from client)
        #   4. If free tier: atomically increment daily counter in Supabase
        #      → enforce 20 msg/day limit server-side
        #
        # Fallback: if Supabase is not configured (local dev), fall back to
        # IP-based in-memory counter so development still works.
        auth_header = request.headers.get('Authorization', '')
        jwt_token   = auth_header[7:] if auth_header.startswith('Bearer ') else ''

        verified_user = _verify_supabase_jwt(jwt_token) if jwt_token else None
        verified_user_id = verified_user.get('id') if verified_user else None

        if verified_user_id:
            # Authenticated user — look up real tier from DB
            server_tier = _get_user_tier_from_db(verified_user_id)
        else:
            # No valid JWT → treat as free/guest regardless of what client sent
            server_tier = 'free'
            # Use IP as fallback identifier for unauthenticated requests
            verified_user_id = f'ip:{get_remote_address()}'

        # Prune old in-memory keys to prevent unbounded growth
        if len(_free_tier_counters) > 50000:
            # FIX: use timedelta for correct date arithmetic (replace(day=...) breaks at month boundaries)
            from datetime import timedelta
            _old_day = (datetime.utcnow() - timedelta(days=2)).strftime('%Y-%m-%d')
            for k in list(_free_tier_counters.keys()):
                if _old_day in k:
                    del _free_tier_counters[k]

        if server_tier == 'free':
            day_key = datetime.utcnow().strftime('%Y-%m-%d')
            _count  = _get_and_increment_daily_count(verified_user_id, day_key)
            if _count > FREE_TIER_DAILY_LIMIT:
                return jsonify({'success': False, 'error': 'Daily message limit reached. Upgrade to Pro for unlimited messages.'}), 429

        # Keep user_tier for any downstream logging / prompt hints
        user_tier = server_tier

        # Parse injected token flags from legacy frontend path
        token_flags = []
        if question.startswith('['):
            tokens = re.findall(r'\[([A-Z_]+)\]', question)
            for tok in tokens:
                token_flags.append(tok)
                question = question.replace(f'[{tok}]', '', 1)
            question = question.strip()

        if 'WEB_SEARCH_ENABLED'  in token_flags: web_search    = True
        if 'THINKING_MODE'       in token_flags: thinking_mode = 'thinking'
        if 'DEEP_THINKING_MODE'  in token_flags: thinking_mode = 'deep'

        logger.info(f"[{mode.upper()}] Q: {question[:80]} | complexity: {complexity}")

        # Select model
        selected_model = MODEL
        if thinking_mode == 'deep':
            selected_model = os.environ.get('DEEP_MODEL', 'deepseek/deepseek-r1:free')
        elif thinking_mode == 'thinking':
            selected_model = os.environ.get('THINK_MODEL', 'deepseek/deepseek-r1-distill-llama-70b:free')

        # Use per-book cached index (no global state race condition)
        searcher = get_book_index(book_id)

        # NEW: subject-agnostic — searches whenever a book is loaded,
        # regardless of whether the question sounds like "chemistry".
        # Works correctly for nursing, biology, physics, etc.
        use_textbook = should_search_textbook(question, chunks_loaded=bool(searcher.chunks))
        logger.info(f"Search textbook: {use_textbook} | book: {book_id}")

        if use_textbook:
            context, similarity, is_relevant, source, all_sources = searcher.smart_search(question, top_k=5)
            logger.debug(f"Score: {similarity:.4f} | Relevant: {is_relevant}")
        else:
            context, similarity, is_relevant, source, all_sources = "", 0.0, False, None, []
            logger.info("Chit-chat / no book loaded")

        # ── Shared prompt helpers ─────────────────────────────────────────

        complexity_levels = {
            1:  "Explain in the simplest possible terms, like to a curious 10-year-old. Use everyday analogies only.",
            2:  "Explain simply for a beginner with no background in this subject. Avoid jargon entirely.",
            3:  "Explain clearly for a middle-school or early high school student. Introduce basic terms gently.",
            4:  "Explain for a high school student. Use standard vocabulary with brief definitions.",
            5:  "Balanced explanation with proper terminology, suitable for an advanced high school or introductory university student.",
            6:  "Detailed explanation for a first-year university student.",
            7:  "University-level explanation. Include relevant equations, mechanisms, and quantitative reasoning where applicable.",
            8:  "Advanced undergraduate level. Use rigorous terminology, derive relationships, and discuss exceptions.",
            9:  "Graduate-level depth. Include theoretical underpinnings and nuanced discussion where relevant.",
            10: "Expert/research level. Provide a comprehensive, highly technical explanation with full mathematical or clinical treatment."
        }
        complexity_instruction = complexity_levels[complexity]

        ctx_block = f"TEXTBOOK CONTEXT (from {BOOK_LIBRARY.get(book_id, {}).get('name', 'textbook')}):\n{context}\n\n" if is_relevant else ""

        # Build user memory block
        memory_block = ""
        if user_memory and user_memory.strip():
            memory_block = f"\n\nUSER PROFILE (remember this about the student):\n{user_memory.strip()}"

        # Book-aware system prompt — works for chemistry, nursing, biology, physics, etc.
        book_info   = BOOK_LIBRARY.get(book_id, {})
        book_name   = book_info.get('name', 'the textbook')
        book_author = book_info.get('author', '')
        book_label  = f"{book_name} by {book_author}" if book_author else book_name

        # Only inject LaTeX instruction for science/math books where equations appear.
        # For a nursing or law book it's irrelevant noise in the prompt.
        EQUATION_SUBJECTS = {'chemistry', 'physics', 'biochemistry', 'mathematics', 'engineering'}
        needs_latex = any(s in book_name.lower() for s in EQUATION_SUBJECTS)
        latex_instruction = (
            "Use LaTeX notation for all equations and formulas. "
            "Inline math: $...$ — Display math: $$...$$ — "
            r"Example: $$K_{eq} = \frac{[C]^c[D]^d}{[A]^a[B]^b}$$"
        ) if needs_latex else (
            "Use plain text for any formulas or technical notation."
        )

        if is_relevant:
            base_system = (
                f"You are an expert tutor for {book_label}. "
                f"Answer based strictly on the provided textbook context and cite page numbers using: 📖 Page N. "
                f"{latex_instruction}{memory_block}"
            )
        else:
            base_system = (
                f"You are a knowledgeable tutor. Answer the student's question helpfully and clearly. "
                f"{latex_instruction}{memory_block}"
            )

        # ── MODE: EXAM ────────────────────────────────────────────────────

        if mode == 'exam':

            # For exam mode we need substantially more context than study mode.
            # 5 chunks (~5 pages) is not enough raw material to write 10
            # genuinely book-grounded questions — the AI fills the gaps with
            # general knowledge. Scale context window up with complexity:
            #   levels 1-4  →  8 chunks  (recall questions)
            #   levels 5-7  → 12 chunks  (application + multi-step)
            #   levels 8-10 → 20 chunks  (book-specific numerical / derivation)
            if complexity <= 4:
                exam_top_k = 8
            elif complexity <= 7:
                exam_top_k = 12
            else:
                exam_top_k = 20

            if use_textbook and searcher.chunks:
                exam_context, exam_similarity, exam_relevant, source, all_sources = \
                    searcher.smart_search(question, top_k=exam_top_k)
                is_relevant = exam_relevant
                similarity  = exam_similarity
            else:
                exam_context = context

            # Exam-specific complexity instructions.
            # These describe question *difficulty and style*, NOT how to explain an answer.
            exam_complexity_levels = {
                1:  ("Write simple recognition questions testing basic vocabulary and definitions. "
                     "Options should be obviously distinct. No calculations required."),
                2:  ("Write recall questions about names, definitions, and basic facts. "
                     "One clearly correct answer, three clearly wrong distractors."),
                3:  ("Write questions where students identify the correct term, formula, or "
                     "simple concept from four options."),
                4:  ("Write straightforward application questions. Include 1-2 questions "
                     "requiring a simple one-step calculation or formula substitution."),
                5:  ("Write mixed recall and application questions. Include 3-4 questions "
                     "requiring multi-step reasoning or formula use. Distractors should be "
                     "plausible misconceptions."),
                6:  ("Write questions requiring understanding of mechanisms and relationships. "
                     "Include 4-5 numerical or equation-based questions. Distractors are "
                     "common student errors."),
                7:  ("Write questions requiring multi-step problem solving. All distractors must "
                     "represent specific calculation errors or conceptual confusions. "
                     "At least 6 questions must involve calculations or derivations."),
                8:  ("Write advanced questions requiring integration of multiple concepts. "
                     "All 10 questions must be calculation or derivation based. "
                     "Use specific numerical values and equations from the textbook context. "
                     "Distractors differ by a common error: sign error, wrong unit, or wrong formula."),
                9:  ("Write graduate-level questions anchored to specific data, equations, or "
                     "worked examples from the textbook context — reference exact values or "
                     "conditions stated on those pages. "
                     "Questions should require derivations, limiting-case analysis, or "
                     "thermodynamic/mechanistic reasoning."),
                10: ("Write research/exam-board level questions using ONLY information "
                     "explicitly present in the textbook pages provided. Every question must:\n"
                     "  - Reference a specific equation, numerical value, figure description, "
                     "or worked example from the context (cite it in the question stem)\n"
                     "  - Require multi-step reasoning: derive, predict, or critically analyse\n"
                     "  - Have distractors that differ by exactly one conceptual or arithmetic error\n"
                     "  - Mirror the style of end-of-chapter problems in university textbooks\n"
                     "Do NOT use any fact, value, or equation not present in the provided pages."),
            }
            exam_complexity_instruction = exam_complexity_levels[complexity]

            exam_ctx_block = (
                f"TEXTBOOK PAGES (base ALL questions on this content only):\n"
                f"{exam_context}\n\n"
            ) if is_relevant else ""

            source_constraint = (
                "CRITICAL: Every question must be directly answerable from the textbook pages "
                "above. Do not introduce facts, values, or equations absent from those pages."
            ) if is_relevant else (
                "Generate questions on this topic at the appropriate difficulty level."
            )

            prompt = f"""You are writing an exam for students studying {book_label}.

{exam_ctx_block}TOPIC: {question}

DIFFICULTY — LEVEL {complexity}/10:
{exam_complexity_instruction}

{source_constraint}

Generate exactly 10 multiple-choice questions.

STRICT FORMAT — follow this exactly for every question:
Q1. [Question text]
A) [option]
B) [option]
C) [option]
D) [option]
Answer: [letter]
Explanation: [Explain why the correct answer is right AND why each wrong option is wrong. Cite the page if you used a specific value: 📖 Page N. Use LaTeX for all equations.]

Q2. ...

Rules:
- All 10 questions must be on the topic above
- Only ONE correct answer per question
- Each question must cover a DIFFERENT concept, calculation, or mechanism
- {latex_instruction}
- Do NOT add any text before Q1 or after Q10's explanation"""

            answer = call_ai(prompt, system_prompt=base_system, model=selected_model, history=history)
            questions = _parse_mcq(answer)
            return jsonify({
                'success': True,
                'mode': 'exam',
                'raw': answer,
                'questions': questions,
                'question_count': len(questions),
                'similarity': float(similarity),
                'is_relevant': is_relevant,
                'source': source,
                'sources': all_sources,
                'complexity_used': complexity
            })

        # ── MODE: PRACTICE ────────────────────────────────────────────────

        elif mode == 'practice':
            prompt = f"""You are a problem-solving tutor for {book_label}.

{ctx_block}TOPIC / QUESTION: {question}

Create a step-by-step problem-solving session at COMPLEXITY LEVEL {complexity}/10: {complexity_instruction}

Structure your response like this:
1. PROBLEM STATEMENT — clearly state a concrete problem to solve (numerical or conceptual).
2. GIVEN — list all given values/information.
3. FIND — state what needs to be determined.
4. SOLUTION — solve step by step, showing every calculation.
5. ANSWER — box the final answer clearly.
6. TIP — give one practical exam tip related to this type of problem.

{latex_instruction}"""

            answer = call_ai(prompt, system_prompt=base_system, model=selected_model, history=history)
            return jsonify({
                'success': True,
                'mode': 'practice',
                'answer': answer,
                'similarity': float(similarity),
                'is_relevant': is_relevant,
                'source': source,
                'sources': all_sources,
                'complexity_used': complexity
            })

        # ── MODE: SUMMARY ─────────────────────────────────────────────────

        elif mode == 'summary':
            prompt = f"""You are a tutor creating a study summary for {book_label}.

{ctx_block}TOPIC: {question}

Write a structured summary at COMPLEXITY LEVEL {complexity}/10: {complexity_instruction}

Include these sections:
1. OVERVIEW — 2–3 sentence big-picture explanation.
2. KEY CONCEPTS — the most important ideas, definitions, and principles.
3. IMPORTANT EQUATIONS — all relevant formulas (use LaTeX).
4. COMMON EXAMPLES — 1–2 real-world or textbook examples.
5. THINGS TO REMEMBER — bullet list of must-know facts and common pitfalls.

{latex_instruction}
Keep the summary focused, clear, and easy to review before an exam."""

            answer = call_ai(prompt, system_prompt=base_system, model=selected_model, history=history)
            return jsonify({
                'success': True,
                'mode': 'summary',
                'answer': answer,
                'similarity': float(similarity),
                'is_relevant': is_relevant,
                'source': source,
                'sources': all_sources,
                'complexity_used': complexity
            })

        # ── MODE: STUDY (default) ─────────────────────────────────────────

        else:
            # ── WEB SEARCH mode ───────────────────────────────────────────
            if web_search:
                web_system = (
                    "You are a helpful research assistant. Answer clearly and accurately "
                    "using current web information. Use markdown formatting with headers, "
                    "bullet points, and bold text where it aids clarity. When you reference "
                    "a source, include the full URL in the text so users can visit it."
                )
                answer, web_citations = call_ai_web_search(
                    question, system_prompt=web_system, history=history
                )
                # Fallback: if web search returned an error, use the standard model
                if answer.startswith('Error:') or answer.startswith('Web search error:'):
                    logger.warning(f"Web search failed ({answer[:80]}), falling back to standard model")
                    fallback_prompt = f"STUDENT QUESTION: {question}\n\nAnswer helpfully and clearly."
                    answer = call_ai(fallback_prompt, system_prompt=base_system, model=selected_model, history=history)
                    answer = "*(Web search unavailable — answering from general knowledge)*\n\n" + answer
                    web_citations = []
                return jsonify({
                    'success': True,
                    'mode': 'study',
                    'answer': answer,
                    'web_search': True,
                    'web_citations': web_citations,
                    'similarity': 0.0,
                    'is_relevant': False,
                    'source': None,
                    'sources': [],
                    'complexity_used': complexity
                })

            # ── Normal textbook / general mode ────────────────────────────
            if is_relevant:
                prompt = f"""You are a tutor for {book_label}.

TEXTBOOK CONTEXT (cite pages using 📖 Page N):
{context}

STUDENT QUESTION: {question}

COMPLEXITY LEVEL {complexity}/10: {complexity_instruction}

FORMATTING: {latex_instruction}

Answer based on the textbook context. Be helpful and clear. Cite the page number whenever you reference specific information from the context."""
            else:
                prompt = f"""You are a knowledgeable tutor.

STUDENT QUESTION: {question}

COMPLEXITY LEVEL {complexity}/10: {complexity_instruction}

FORMATTING: {latex_instruction}

Answer helpfully and clearly."""

            answer = call_ai(prompt, system_prompt=base_system, model=selected_model, history=history)
            return jsonify({
                'success': True,
                'mode': 'study',
                'answer': answer,
                'context': context,
                'similarity': float(similarity),
                'is_relevant': is_relevant,
                'source': source,
                'sources': all_sources,   # FIX: frontend uses data.sources for badge row
                'complexity_used': complexity
            })

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


def _parse_mcq(raw_text):
    questions = []
    blocks = re.split(r'\n(?=Q\d+\.)', raw_text.strip())
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        lines = block.splitlines()
        q_obj = {'number': None, 'question': '', 'options': {}, 'answer': '', 'explanation': ''}
        for line in lines:
            line = line.strip()
            m = re.match(r'^Q(\d+)\.\s*(.*)', line)
            if m:
                q_obj['number'] = int(m.group(1))
                q_obj['question'] = m.group(2)
                continue
            m = re.match(r'^([A-D])[).]\s*(.*)', line)
            if m:
                q_obj['options'][m.group(1)] = m.group(2)
                continue
            m = re.match(r'^Answer:\s*(.*)', line, re.IGNORECASE)
            if m:
                q_obj['answer'] = m.group(1).strip()
                continue
            m = re.match(r'^Explanation:\s*(.*)', line, re.IGNORECASE)
            if m:
                q_obj['explanation'] = m.group(1).strip()
        if q_obj['number'] is not None:
            questions.append(q_obj)
    return questions


# ============================================
# /generate-flashcards
# ============================================

@app.route('/generate-flashcards', methods=['POST', 'OPTIONS'])
@limiter.limit('10 per minute; 60 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def generate_flashcards():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        data = request.json
        if not data:
            return jsonify({'success': False, 'error': 'Invalid or missing JSON body'}), 400

        topic   = data.get('topic', 'chemistry').strip()
        count   = min(int(data.get('count', 10)), 20)
        book_id = data.get('bookId', 'zumdahl')

        # ── Cache check: return instantly if already generated ────────────────
        cache_k = _cache_key(book_id, topic, 'flashcards', count)
        cached  = _cache_get(cache_k)
        if cached:
            logger.info(f"⚡ Cache HIT flashcards: {topic} ({book_id})")
            return jsonify({**cached, 'cached': True})
        logger.info(f"🔄 Cache MISS flashcards: {topic} ({book_id})")

        context_block = ""
        searcher = get_book_index(book_id)
        if searcher.chunks:
            context, score, is_relevant, _, _ = searcher.smart_search(topic, top_k=3)
            if is_relevant:
                context_block = f"Use this textbook content as your primary source:\n{context}\n\n"

        prompt = f"""{context_block}Create exactly {count} flashcards about: {topic}

STRICT OUTPUT FORMAT — output ONLY this, no intro text, no numbering prose:
CARD
FRONT: [concise question or term — max 20 words]
BACK: [clear precise answer — max 60 words, use LaTeX for equations: $...$]
END

Repeat the CARD / FRONT / BACK / END block exactly {count} times.
Rules:
- Cover definitions, equations, mechanisms, and key facts
- Each card must be self-contained
- No duplicate questions
- Use $LaTeX$ for all formulas/equations"""

        raw = call_ai(prompt, system_prompt=(
            "You are a chemistry flashcard generator. Output ONLY the CARD blocks in the exact format requested. "
            "No preamble, no extra commentary, no numbering outside the format."
        ), model=MODEL)

        flashcards = []
        blocks = re.split(r'\bCARD\b', raw, flags=re.IGNORECASE)
        for block in blocks:
            block = block.strip()
            if not block:
                continue
            front_match = re.search(r'FRONT:\s*(.+?)(?=BACK:|$)', block, re.IGNORECASE | re.DOTALL)
            back_match  = re.search(r'BACK:\s*(.+?)(?=END|CARD|$)', block, re.IGNORECASE | re.DOTALL)
            if front_match and back_match:
                front = front_match.group(1).strip().rstrip('END').strip()
                back  = back_match.group(1).strip().rstrip('END').strip()
                if front and back:
                    flashcards.append({'front': front, 'back': back})

        if not flashcards:
            # Fallback: Q:/A: format
            for block in re.split(r'\n(?=Q\d*[:.\s])', raw):
                q_match = re.search(r'Q\d*[:.\s]\s*(.+?)(?=A[:.\s]\s|$)', block, re.DOTALL)
                a_match = re.search(r'A[:.\s]\s*(.+)', block, re.DOTALL)
                if q_match and a_match:
                    flashcards.append({'front': q_match.group(1).strip(), 'back': a_match.group(1).strip()})

        if not flashcards:
            return jsonify({'success': False, 'error': 'Failed to parse flashcards', 'raw': raw}), 500

        logger.info(f"Generated {len(flashcards)} flashcards for: {topic}")
        result_payload = {'success': True, 'flashcards': flashcards, 'count': len(flashcards), 'topic': topic}
        _cache_set(cache_k, result_payload)  # FIX: actually store result in cache
        return jsonify(result_payload)

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================
# /upload-document
# ============================================

@app.route('/upload-document', methods=['POST', 'OPTIONS'])
@limiter.limit('10 per minute; 50 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def upload_document():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not file.filename:
            return jsonify({'success': False, 'error': 'Empty filename'}), 400
        safe_name = secure_filename(file.filename)
        if not safe_name or not allowed_file(safe_name):
            return jsonify({'success': False, 'error': 'Unsupported file type. Allowed: PDF, DOCX, PPTX'}), 400
        filename  = safe_name.lower()
        # FIX: use uuid to avoid race condition when two uploads happen at the same second
        import uuid as _uuid
        temp_path = f"/tmp/chunks_{_uuid.uuid4().hex}_{safe_name}"
        file.save(temp_path)
        extracted_slides = []

        try:
            if filename.endswith(('.pptx', '.ppt')):
                prs = Presentation(temp_path)
                for i, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    slide_title = f"Slide {i}"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # shape_type==13 is PICTURE (bug fix: was incorrectly used as title check).
                            # Correct: use placeholder index 0 (PP_PLACEHOLDER.TITLE) for title detection.
                            is_title = (
                                hasattr(shape, "is_placeholder") and shape.is_placeholder and
                                hasattr(shape, "placeholder_format") and shape.placeholder_format is not None and
                                shape.placeholder_format.idx == 0
                            )
                            if is_title:
                                slide_title = text
                            else:
                                slide_texts.append(text)
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = " | ".join(
                                    cell.text.strip() for cell in row.cells if cell.text.strip()
                                )
                                if row_text:
                                    slide_texts.append(row_text)
                    notes_text = ""
                    if slide.has_notes_slide:
                        nf = slide.notes_slide.notes_text_frame
                        if nf:
                            notes_text = nf.text.strip()
                    extracted_slides.append({
                        'slide_number': i, 'title': slide_title,
                        'content': slide_texts, 'notes': notes_text
                    })

            elif filename.endswith('.pdf'):
                with open(temp_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages, 1):
                        text = page.extract_text() or ""
                        extracted_slides.append({
                            'slide_number': i, 'title': f"Page {i}",
                            'content': [text] if text.strip() else [], 'notes': ''
                        })

            elif filename.endswith('.docx'):
                doc = docx.Document(temp_path)
                current_section = {"title": safe_name, "content": [], "slide_number": 1, "notes": ""}
                for para in doc.paragraphs:
                    if not para.text.strip():
                        continue
                    if para.style.name.startswith('Heading'):
                        if current_section["content"]:
                            extracted_slides.append(current_section)
                        current_section = {
                            "title": para.text.strip(), "content": [],
                            "slide_number": len(extracted_slides) + 1, "notes": ""
                        }
                    else:
                        current_section["content"].append(para.text.strip())
                if current_section["content"]:
                    extracted_slides.append(current_section)
                if not extracted_slides:
                    all_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                    extracted_slides.append({
                        'slide_number': 1, 'title': safe_name,
                        'content': [all_text], 'notes': ''
                    })
            else:
                os.remove(temp_path)
                return jsonify({'success': False, 'error': f'Unsupported file type: {filename}'}), 400

            os.remove(temp_path)

            total_text = " ".join(" ".join(s.get('content', [])) for s in extracted_slides)
            if len(total_text.strip()) < 30:
                return jsonify({
                    'success': False,
                    'error': 'Could not extract readable text. The file may be scanned/image-based or empty.'
                }), 400

            return jsonify({
                'success': True,
                'slides': extracted_slides,
                'total_slides': len(extracted_slides),
                'filename': safe_name
            })

        except Exception as e:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            raise e

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================
# /generate-study-materials
# ============================================

@app.route('/generate-study-materials', methods=['POST', 'OPTIONS'])
@limiter.limit('5 per minute; 20 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def generate_study_materials():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        data          = request.json
        slides        = data.get('slides', [])
        material_type = data.get('type', 'notes')

        # ── Cache check ───────────────────────────────────────────────────────
        import hashlib as _hl
        _sm_hash = _hl.md5(str(slides).encode()).hexdigest()[:16]
        _sm_cache_k = _cache_key('doc', _sm_hash, material_type, 0)
        _sm_cached = _cache_get(_sm_cache_k)
        if _sm_cached:
            logger.info(f"⚡ Cache HIT study-materials: {material_type}")
            return jsonify({**_sm_cached, 'cached': True})

        if not slides:
            return jsonify({'success': False, 'error': 'No slide content provided'}), 400

        full_content = ""
        for slide in slides:
            title = slide.get('title', f"Slide {slide.get('slide_number', '?')}")
            content_lines = [l for l in slide.get('content', []) if l.strip()]
            notes = slide.get('notes', '').strip()
            if not content_lines and not notes:
                continue
            full_content += f"\n\n=== {title} ===\n"
            full_content += "\n".join(content_lines)
            if notes:
                full_content += f"\n[Speaker Notes: {notes}]"

        if not full_content.strip():
            return jsonify({'success': False, 'error': 'No readable content found in slides'}), 400

        # FIX: increased cap from 12,000 to 24,000 chars to cover more slides.
        # Also added a visible truncation warning in the prompt so the AI knows
        # it may be working with partial content (was silently truncating before).
        char_limit = 24000
        content_for_ai = full_content.strip()[:char_limit]
        truncation_note = ""
        if len(full_content) > char_limit:
            truncation_note = (
                f"\n\n[NOTE: The slide content was trimmed to fit the model's context window. "
                f"The above represents the first {char_limit} characters. "
                f"Focus your output on what is provided — do not guess at omitted content.]"
            )
            content_for_ai += truncation_note

        if material_type == 'notes':
            prompt = f"""You are a meticulous academic note-taker. Below is the EXACT text extracted from a student's lecture slides.

YOUR TASK: Write comprehensive, well-organized study notes based STRICTLY on the content below.

RULES:
- Cover EVERY topic, concept, definition, formula, and example that appears in the slides
- Do NOT add information that is not in the slides
- Do NOT skip any slide or section
- Use clear headings that match the slide titles
- Preserve all technical terms, formulas, and specific details exactly as they appear
- If a slide has bullet points, expand them into full explanatory sentences
- Format: use markdown with ## for main sections, ### for subsections, and bullet points for lists

SLIDE CONTENT:
{content_for_ai}

Now write the complete study notes:"""

        elif material_type == 'reviewer':
            prompt = f"""You are a senior academic exam coach who has written board review books. Below is the EXACT text from a student's lecture slides.

YOUR TASK: Create a comprehensive, high-yield exam reviewer from this content. Be thorough and precise — students depend on this to pass their exams.

OUTPUT STRUCTURE (use these exact markdown headings):

## 📌 Topic Overview
Write 3–5 sentences summarizing the core theme of this material and what students are expected to master.

## 🔑 High-Yield Concepts
List every testable concept, definition, mechanism, and principle found in the slides. For each:
- **Term/Concept**: Clear, exam-ready definition or explanation
- Flag with ⚠️ anything that appears repeatedly or is emphasized in the slides (likely to be on the exam)

## 🧮 Formulas & Equations
List every formula, equation, or quantitative relationship in the slides. For each:
- Write it clearly with variable definitions
- Add a one-line note on when/how to apply it

## 📊 Key Values, Constants & Comparisons
Extract all specific numbers, thresholds, units, classifications, or comparative data (e.g., temperature ranges, normal values, conversion factors, rankings).

## 🧠 Mnemonics & Memory Tips
Create 3–6 original mnemonics or memory hooks for the hardest-to-remember facts from the slides.

## ❓ Practice Questions (10 questions)
Write 10 exam-style questions directly from the slide content. Mix question types:
- 5 multiple choice (A/B/C/D with answer and explanation)
- 3 short-answer (with model answer)
- 2 "explain why" or "compare and contrast" style

For each MCQ:
Q: [question]
A) | B) | C) | D)
✅ Answer: [letter] — [1-sentence explanation referencing the slide content]

## 🚨 Common Mistakes to Avoid
List 3–5 common misconceptions or errors students make on this topic, based on what the slides emphasize.

## ⚡ Last-Minute Cheat Sheet
A condensed bullet list of the 10–15 most important facts to review the night before the exam.

RULES:
- Base EVERYTHING strictly on the provided slide content
- Do NOT fabricate, guess, or add external information
- If slides mention a specific number, value, or name — include it exactly
- Be as detailed as the content allows — do NOT be generic
- Use markdown formatting throughout

SLIDE CONTENT:
{content_for_ai}

Write the complete exam reviewer:"""

        elif material_type == 'flashcards':
            prompt = f"""You are a flashcard creator. Below is the EXACT text from lecture slides.

YOUR TASK: Create 15-20 flashcards based STRICTLY on this content.

FORMAT for each card:
CARD [N]
Q: [Question about a specific concept, term, or fact from the slides]
A: [Precise answer drawn directly from the slide content]

RULES:
- Every question must be answerable using the slide content provided
- Cover key terms, definitions, processes, formulas, and important facts
- Do NOT create questions about topics not in the slides

SLIDE CONTENT:
{content_for_ai}

Create the flashcards:"""

        elif material_type == 'summary':
            prompt = f"""You are a study guide writer. Below is the EXACT text from lecture slides.

YOUR TASK: Create a concise one-page summary sheet based STRICTLY on this content.

Include:
1. MAIN TOPIC & OVERVIEW (2-3 sentences)
2. KEY CONCEPTS & DEFINITIONS (from the slides only)
3. IMPORTANT FORMULAS/EQUATIONS (if any appear in the slides)
4. CRITICAL FACTS TO REMEMBER (the most testable points from the slides)

RULES:
- Only include what is in the slides
- Be concise but complete

SLIDE CONTENT:
{content_for_ai}

Write the summary sheet:"""

        elif material_type == 'quiz':
            prompt = f"""You are a quiz generator. Below is the EXACT text from lecture slides.

YOUR TASK: Generate exactly 10 multiple-choice questions based STRICTLY on this content.

STRICT FORMAT:
Q1. [Question text based on the slide content]
A) [option]
B) [option]
C) [option]
D) [option]
Answer: [correct letter]
Explanation: [brief explanation referencing the slide content]

RULES:
- All 10 questions must be directly answerable from the slide content provided
- Do NOT ask about topics not covered in the slides
- Only ONE correct answer per question

SLIDE CONTENT:
{content_for_ai}

Generate the quiz:"""

        elif material_type == 'all':
            prompt = f"""You are a comprehensive study material generator. Below is the EXACT text from lecture slides.

YOUR TASK: Generate ALL of the following based STRICTLY on this content:

1. STUDY NOTES - Comprehensive notes covering every topic in the slides
2. KEY FLASHCARDS - 10 Q&A cards for the most important concepts
3. SUMMARY SHEET - A concise overview of the main points
4. 5 PRACTICE QUESTIONS - Multiple choice questions from the slide content

RULES:
- Base EVERYTHING strictly on the provided slide content
- Do NOT add external information
- Label each section clearly

SLIDE CONTENT:
{content_for_ai}

Generate all study materials:"""

        else:
            prompt = f"""You are an academic assistant. Below is the EXACT text from lecture slides.

Create detailed {material_type} based STRICTLY on this content. Do not add information not present in the slides.

SLIDE CONTENT:
{content_for_ai}"""

        result = call_ai(prompt, system_prompt=(
            "You are an expert academic assistant and exam coach who creates precise, high-yield study materials. "
            "You ONLY use information from the provided source content — never fabricate or hallucinate. "
            "For exam reviewers, you think like a professor writing the exam: you identify what is most testable, "
            "what students most commonly get wrong, and what concepts are foundational vs peripheral. "
            "You write in a clear, structured format that students can scan quickly under exam pressure."
        ), max_tokens_override=8000)

        sm_payload = {'success': True, 'materials': {material_type: result}}
        _cache_set(_sm_cache_k, sm_payload)  # FIX: actually store result in cache
        return jsonify(sm_payload)

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================
# /generate-quiz
# ============================================

@app.route('/generate-quiz', methods=['POST', 'OPTIONS'])
@limiter.limit('10 per minute; 60 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def generate_quiz():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        data       = request.json or {}
        slides     = data.get('slides', [])
        count      = max(5, min(50, int(data.get('count', 10))))
        difficulty = data.get('difficulty', 'medium').lower().strip()
        quiz_mode  = data.get('mode', 'standard').lower().strip()
        existing_questions = data.get('existingQuestions', [])

        if not slides:
            return jsonify({'success': False, 'error': 'No slide content provided'}), 400

        full_content = ""
        for slide in slides:
            title = slide.get('title', f"Slide {slide.get('slide_number','?')}")
            content_lines = [l for l in slide.get('content', []) if l.strip()]
            notes = slide.get('notes', '').strip()
            if not content_lines and not notes:
                continue
            full_content += f"\n\n=== {title} ===\n" + "\n".join(content_lines)
            if notes:
                full_content += f"\n[Notes: {notes}]"

        content_for_ai = full_content.strip()[:24000]  # FIX: was 12000
        if not content_for_ai:
            return jsonify({'success': False, 'error': 'No readable content in slides'}), 400

        difficulty_instructions = {
            'easy': (
                "DIFFICULTY: EASY\n"
                "- Test direct recall of terms, definitions, and basic facts stated explicitly in the slides\n"
                "- Questions should be clear and unambiguous — one obviously correct answer\n"
                "- Wrong options (distractors) should be clearly wrong but come from related topics in the slides\n"
                "- Bloom's Taxonomy: Knowledge/Remember level\n"
                "- Example question style: 'What is the definition of...', 'Which of the following is...', 'According to the slides, ...'"
            ),
            'medium': (
                "DIFFICULTY: MEDIUM\n"
                "- Test understanding and application of concepts from the slides\n"
                "- Questions should require the student to understand WHY, not just WHAT\n"
                "- Distractors should be plausible — common misunderstandings, closely related concepts, or partially correct statements\n"
                "- Include at least 2 application questions ('A student does X, what happens?')\n"
                "- Bloom's Taxonomy: Understand/Apply level\n"
                "- Example question style: 'Which best explains why...', 'If X occurs, what is the result...', 'What is the relationship between...'"
            ),
            'hard': (
                "DIFFICULTY: HARD\n"
                "- Test analysis, synthesis, and evaluation — the highest Bloom's levels\n"
                "- Questions should require integrating multiple concepts from across the slides\n"
                "- Distractors must be very close to the correct answer — they should be statements that are true in a DIFFERENT context or partially correct\n"
                "- Include at least 3 'EXCEPT' type questions (e.g., 'All of the following are true EXCEPT...')\n"
                "- Include comparison questions and 'what would happen if...' type scenarios\n"
                "- No question should be answerable by simple recall alone\n"
                "- Bloom's Taxonomy: Analyze/Evaluate/Create level"
            ),
        }
        diff_text = difficulty_instructions.get(difficulty, difficulty_instructions['medium'])

        no_repeat_block = ""
        if existing_questions:
            sample = existing_questions[:30]
            no_repeat_block = (
                "\n\nIMPORTANT — DO NOT REPEAT: The following question topics have already been used. "
                "Generate completely NEW questions on DIFFERENT aspects of the content:\n"
                + "\n".join(f"- {q}" for q in sample) + "\n"
            )

        if quiz_mode == 'situational':
            mode_instruction = (
                "\nQUIZ MODE: SITUATIONAL\n"
                "- Every question MUST open with a 2–4 sentence real-world scenario or case study\n"
                "- Scenarios should involve: a lab experiment, a patient/student/professional encountering this concept, "
                "a real-world application, or an observed phenomenon\n"
                "- The question should ask the student to diagnose, explain, predict, or decide based on the scenario\n"
                "- All 4 options must be plausible within the context of the scenario\n"
                "- The scenario should be rich enough that the wrong answer choices feel genuinely tempting\n"
                "- Never just restate a fact — always embed it in a realistic context"
            )
        else:
            mode_instruction = ""

        # Smart coverage: ensure questions span the full content, not just early slides
        slide_count = len([s for s in slides if any(s.get('content', []))])
        coverage_note = ""
        if slide_count > 5:
            coverage_note = (
                f"\n\nCOVERAGE REQUIREMENT: The slide deck has {slide_count} slides with content. "
                f"Spread your {count} questions across the ENTIRE deck — do not concentrate them only on the first few slides. "
                "Aim to cover every major section."
            )

        prompt = f"""You are an expert exam writer who creates board-quality multiple choice questions for university students.
Below is the EXACT text from a student's lecture slides.

YOUR TASK: Generate exactly {count} multiple-choice questions based STRICTLY on this content.

{diff_text}{mode_instruction}{coverage_note}{no_repeat_block}

STRICT OUTPUT FORMAT — follow this exactly for every question, no deviations:
Q1. [Question text — write a complete, grammatically correct question]
A) [option — complete sentence or phrase, not just a word]
B) [option]
C) [option]
D) [option]
Answer: [single letter: A, B, C, or D]
Explanation: [3–5 sentences: (1) why the correct answer is right with specific evidence from the slides, (2) why the most tempting wrong answer is wrong, (3) a memory tip or key insight to help the student remember this for the exam]

QUALITY RULES:
- Every question must test a DIFFERENT concept — no two questions on the same fact
- All {count} questions must be answerable from the slide content below — no outside knowledge
- Only ONE correct answer per question — the other three must be clearly wrong (but tempting)
- Questions must be SPECIFIC — never ask vague questions like "Which is important?" 
- Options must all be the same grammatical form and similar length
- Never start questions with "According to the slides" — write naturally
- Do NOT add headers, numbering schemes, or commentary — just Q1 through Q{count} in sequence
- Start immediately with Q1 — no preamble

SLIDE CONTENT:
{content_for_ai}

Generate the quiz:"""

        raw = call_ai(prompt, system_prompt=(
            "You are an expert exam writer with 20 years of experience creating board-level multiple choice questions. "
            "You write questions that genuinely test understanding, not just memorization. "
            "Your distractors are carefully crafted to catch common misconceptions. "
            "You strictly use only information from the provided source material — never add external knowledge. "
            "You always follow the exact output format with no deviations. "
            "You write thorough, educational explanations that help students learn from both correct and incorrect answers."
        ), max_tokens_override=12000)

        questions = _parse_mcq(raw)
        if not questions:
            return jsonify({'success': False, 'error': 'Could not parse quiz output. Try again.', 'raw': raw}), 500

        return jsonify({
            'success': True,
            'questions': questions,
            'count': len(questions),
            'difficulty': difficulty,
            'raw': raw
        })

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================
# /ask-image  — FIX: was missing @app.route decorator (complete 404)
# ============================================

@app.route('/ask-image', methods=['POST', 'OPTIONS'])
@limiter.limit('10 per minute; 40 per hour', exempt_when=lambda: request.method == 'OPTIONS')
def ask_image():
    if request.method == 'OPTIONS':
        return jsonify({'ok': True})
    try:
        data       = request.json
        question   = data.get('question', 'Describe what you see and explain any chemistry concepts visible.')
        image_b64  = data.get('image_b64', '')
        image_type = data.get('image_type', 'image/jpeg')
        complexity = data.get('complexity', 5)
        thinking   = data.get('thinking', None)

        if not image_b64:
            return jsonify({'success': False, 'error': 'No image data provided'}), 400

        vision_model = os.environ.get('VISION_MODEL', 'nvidia/nemotron-nano-12b-v2-vl:free')

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type":  "application/json",
            "HTTP-Referer":  "https://chunks-ai.vercel.app",
            "X-Title":       "Chunks Chemistry"
        }

        complexity_levels = {
            1: "very simple terms a child can understand",
            2: "simple terms for a beginner",
            3: "middle-school level",
            4: "high school level",
            5: "AP/IB chemistry level",
            6: "first-year university level",
            7: "university level with equations",
            8: "advanced undergraduate",
            9: "graduate level",
            10: "expert/research level"
        }
        level_desc = complexity_levels.get(max(1, min(10, int(complexity))), "university level")

        system_prompt = (
            f"You are an expert chemistry tutor with vision capabilities. "
            f"Analyze the image carefully and explain any chemistry concepts, "
            f"diagrams, equations, molecules, lab setups, or periodic table elements visible. "
            f"Explain at {level_desc}. "
            f"Use LaTeX for equations: inline $...$ and display $$...$$."
        )

        payload = {
            "model": vision_model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:{image_type};base64,{image_b64}"}},
                        {"type": "text", "text": question}
                    ]
                }
            ],
            "temperature": 0.15,  # FIX: was 0.4
            "max_tokens": 2000
        }

        logger.info(f"Vision model: {vision_model}")
        response = _session.post(OPENROUTER_URL, headers=headers, json=payload, timeout=55)

        if response.status_code == 200:
            answer = response.json()['choices'][0]['message']['content']
            return jsonify({'success': True, 'answer': answer, 'model': vision_model})
        else:
            err_detail = response.text[:400]
            logger.error(f"Vision API error {response.status_code}: {err_detail}")
            return jsonify({
                'success': False,
                'error': f'Vision API error {response.status_code}',
                'detail': err_detail
            }), 500

    except Exception as e:
        logger.exception("Unhandled error")
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================
# ADMIN: OPENROUTER CREDIT USAGE
# ============================================

@app.route('/api/admin/openrouter-credits', methods=['GET', 'OPTIONS'])
def openrouter_credits():
    """Proxy endpoint: returns OpenRouter account credit/usage info to authorized admins.
    Requires a valid Supabase JWT from an admin account passed as Bearer token.
    """
    if request.method == 'OPTIONS':
        return jsonify({}), 200

    # Validate caller is an authenticated admin via Supabase JWT
    auth_header = request.headers.get('Authorization', '')
    if not auth_header.startswith('Bearer '):
        return jsonify({'success': False, 'error': 'Unauthorized'}), 401

    # FIX: Actually verify the JWT and check admin role (was only checking prefix before)
    jwt_token = auth_header[7:]
    verified_admin = _verify_supabase_jwt(jwt_token)
    if not verified_admin:
        logger.warning('Admin endpoint: invalid or expired JWT')
        return jsonify({'success': False, 'error': 'Unauthorized — invalid token'}), 401

    admin_user_id = verified_admin.get('id', '')
    if SUPABASE_URL and SUPABASE_SERVICE_KEY and admin_user_id:
        try:
            resp = _session.get(
                f"{SUPABASE_URL}/rest/v1/users",
                params={"id": f"eq.{admin_user_id}", "select": "role,tier"},
                headers={
                    "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
                    "apikey": SUPABASE_SERVICE_KEY,
                },
                timeout=5
            )
            if resp.status_code == 200:
                rows = resp.json()
                user_role = rows[0].get('role', '') if rows else ''
                if user_role not in ('admin', 'superadmin'):
                    logger.warning(f'Admin endpoint: non-admin user {admin_user_id} attempted access')
                    return jsonify({'success': False, 'error': 'Forbidden — admin role required'}), 403
        except Exception as e:
            logger.warning(f'Admin role check failed: {e}')
            return jsonify({'success': False, 'error': 'Could not verify admin role'}), 500

    try:
        # Fetch key info from OpenRouter
        key_resp = _session.get(
            'https://openrouter.ai/api/v1/auth/key',
            headers={'Authorization': f'Bearer {OPENROUTER_API_KEY}'},
            timeout=10
        )
        if key_resp.status_code != 200:
            return jsonify({'success': False, 'error': f'OpenRouter returned {key_resp.status_code}'}), 502

        key_data = key_resp.json().get('data', {})
        logger.info(f"OpenRouter key_data raw: {key_data}")

        # key_data.usage is the authoritative total spend in USD credits
        key_usage = float(key_data.get('usage', 0) or 0)
        key_limit = key_data.get('limit')  # None means unlimited

        # Fetch generation history for per-model breakdown
        gen_resp = _session.get(
            'https://openrouter.ai/api/v1/generation',
            headers={'Authorization': f'Bearer {OPENROUTER_API_KEY}'},
            params={'limit': 500, 'offset': 0},
            timeout=15
        )
        logger.info(f"OR generation status: {gen_resp.status_code}, body: {gen_resp.text[:300]}")

        total_tokens = 0
        total_requests = 0
        model_breakdown = {}

        if gen_resp.status_code == 200:
            raw = gen_resp.json()
            gen_data = raw.get('data', raw) if isinstance(raw, dict) else raw
            if isinstance(gen_data, list):
                total_requests = len(gen_data)
                for g in gen_data:
                    cost   = float(g.get('total_cost', 0) or 0)
                    tokens = (int(g.get('tokens_prompt', 0) or 0) +
                              int(g.get('tokens_completion', 0) or 0))
                    model  = g.get('model', 'unknown')
                    total_tokens += tokens
                    if model not in model_breakdown:
                        model_breakdown[model] = {'cost': 0.0, 'tokens': 0, 'requests': 0}
                    model_breakdown[model]['cost']     += cost
                    model_breakdown[model]['tokens']   += tokens
                    model_breakdown[model]['requests'] += 1

        # Use key-level usage as the authoritative total spend
        total_cost_usd = key_usage if key_usage > 0 else sum(
            v['cost'] for v in model_breakdown.values()
        )

        remaining = None
        if key_limit is not None:
            remaining = round(float(key_limit) - key_usage, 6)

        return jsonify({
            'success': True,
            'key_info': {
                'label': key_data.get('label', ''),
                'limit': key_limit,
                'limit_remaining': remaining,
                'usage': key_usage,
                'is_free_tier': key_data.get('is_free_tier', False),
                'rate_limit': key_data.get('rate_limit', {}),
            },
            'usage_summary': {
                'total_cost_usd': round(total_cost_usd, 6),
                'total_tokens': total_tokens,
                'total_requests': total_requests,
            },
            'model_breakdown': model_breakdown,
            '_debug': {
                'generation_status': gen_resp.status_code,
                'key_usage_raw': key_usage,
            }
        })

    except Exception as e:
        logger.exception('openrouter_credits error')
        return jsonify({'success': False, 'error': str(e)}), 500

# ── PAEV — Prerequisite-Aware Epistemic Verification ─────────────────────────
from paev_routes import register_paev
register_paev(app)

from progress_routes import register_progress
register_progress(app)
# ============================================
# START SERVER
# ============================================

if __name__ == '__main__':
    logger.info("=" * 60)
    logger.info("🧪 CHUNKS CHEMISTRY - PRODUCTION SERVER")
    logger.info("=" * 60)
    logger.info(f"Mode: {'PRODUCTION' if PRODUCTION else 'DEVELOPMENT'}")
    logger.info(f"Books: {len(BOOK_LIBRARY)}")
    logger.info(f"R2: {'Configured' if R2_BUCKET_URL != 'https://pub-xxxxx.r2.dev' else 'NOT CONFIGURED'}")
    logger.info(f"API: {'Configured' if OPENROUTER_API_KEY != 'your-key-here' else 'NOT CONFIGURED'}")
    logger.info(f"Port: {PORT}")
    # Dev only — use gunicorn in production (Railway uses Procfile)
    app.run(host='0.0.0.0', port=PORT, debug=False)
